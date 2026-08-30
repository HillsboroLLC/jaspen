import csv
import io
import re
from datetime import date, datetime, timezone

from flask import Blueprint, jsonify, request, send_file
from flask_jwt_extended import get_jwt_identity, jwt_required

from app.billing_config import PLAN_RANK, to_public_plan
from app.models import Organization, User
from app.orgs import active_membership_for_user, resolve_active_org_for_user

from .ai_agent import _normalize_analysis_history
from .reports import _markdown_to_pdf_bytes, _safe_text
from .sessions import load_user_sessions
from .strategy import _load_scenarios, _resolve_thread_wbs, _scorecard_snapshot_state

export_bp = Blueprint("export", __name__)

PPTX_MIMETYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
CSV_MIMETYPE = "text/csv; charset=utf-8"
PDF_MIMETYPE = "application/pdf"
XLSX_MIMETYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DOCX_MIMETYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MARKDOWN_MIMETYPE = "text/markdown; charset=utf-8"


def _iso_now():
    return datetime.utcnow().isoformat()


def _resolve_thread_session(sessions, thread_id):
    tid = str(thread_id or "").strip()
    if not tid or not isinstance(sessions, dict):
        return None
    if tid in sessions and isinstance(sessions.get(tid), dict):
        return sessions.get(tid)
    for candidate in sessions.values():
        if not isinstance(candidate, dict):
            continue
        if str(candidate.get("session_id") or "").strip() == tid:
            return candidate
    return None


def _resolve_export_context(user, session):
    plan_key = to_public_plan(getattr(user, "subscription_plan", None))
    org = None
    membership = None
    session_org_id = str((session or {}).get("organization_id") or "").strip()

    if session_org_id:
        org = Organization.query.filter_by(id=session_org_id).first()
        if org is not None:
            membership = active_membership_for_user(session_org_id, user.id)
            if membership:
                plan_key = to_public_plan(org.plan_key)
            else:
                org = None
    elif isinstance(user, User):
        active_org, active_membership = resolve_active_org_for_user(user)
        if active_org and active_membership:
            org = active_org
            membership = active_membership

    return plan_key, org, membership


def _require_export_plan(plan_key, export_type):
    # Customer-facing artifact downloads are available on every signed-in plan.
    # Keep the legacy scorecard Word/Excel endpoints gated for backwards
    # compatibility, but the supported UI only offers PDF and PowerPoint.
    if export_type in {"scorecard_pdf", "scorecard_pptx", "csv", "wbs_xlsx"}:
        return None

    required_plan = "team" if export_type == "pptx" else "essential"
    if PLAN_RANK.get(plan_key, 0) >= PLAN_RANK.get(required_plan, 0):
        return None

    label = {
        "pdf": "PDF export",
        "pptx": "PowerPoint export",
        "scorecard_pdf": "Scorecard PDF export",
        "scorecard_pptx": "Scorecard PowerPoint export",
        "csv": "WBS CSV export",
        "wbs_xlsx": "Execution plan Excel export",
        "xlsx": "Excel export",
        "docx": "Word export",
    }.get(export_type, "Export")
    return (
        jsonify(
            {
                "error": f"{label} is not available on your current plan.",
                "code": "export_plan_required",
                "plan_key": plan_key,
                "required_plan": required_plan,
                "export_type": export_type,
            }
        ),
        403,
    )


def _format_label(key):
    token = str(key or "").replace("_", " ").strip()
    return token.title() if token else "Item"


def _scorecard_record_for_export(session, thread_id, scorecard_id=None, user_id=None):
    analyses = _normalize_analysis_history(session, thread_id)
    scorecard_id = str(scorecard_id or "").strip()
    selected = None

    if user_id is not None:
        from app.scenarios_store import load_scenarios_data
        from app.scorecards import collect_peer_scorecards
        thread_data = (load_scenarios_data(user_id) or {}).get(thread_id) or {}
        peers = collect_peer_scorecards(
            user_id,
            thread_id,
            legacy_session=session,
            legacy_thread_data=thread_data,
        )
        chosen = None
        if scorecard_id:
            chosen = next((
                item for item in peers
                if str(item.get('id') or item.get('analysis_id') or '') == scorecard_id
            ), None)
            if chosen is None:
                return None, (
                    jsonify({"error": "Scorecard not found in this thread.", "code": "scorecard_not_found"}),
                    404,
                )
        elif peers:
            chosen = peers[0]
        if isinstance(chosen, dict):
            selected = {
                'analysis_id': str(chosen.get('id') or chosen.get('analysis_id') or thread_id),
                'created_at': chosen.get('createdAt') or chosen.get('timestamp'),
                'result': chosen,
            }

    if scorecard_id and selected is None:
        for item in analyses:
            if not isinstance(item, dict):
                continue
            result = item.get("result") if isinstance(item.get("result"), dict) else {}
            candidates = {
                str(item.get("analysis_id") or "").strip(),
                str(item.get("id") or "").strip(),
                str(result.get("analysis_id") or "").strip(),
                str(result.get("id") or "").strip(),
                str(result.get("scenario_id") or "").strip(),
            }
            candidates.discard("")
            if scorecard_id in candidates:
                selected = item
                break
        if selected is None:
            return None, (
                jsonify({"error": "Scorecard not found in this thread.", "code": "scorecard_not_found"}),
                404,
            )

    if selected is None:
        result_blob = session.get("result") if isinstance(session.get("result"), dict) else {}
        if result_blob:
            selected = {
                "analysis_id": str(
                    result_blob.get("analysis_id")
                    or result_blob.get("id")
                    or session.get("session_id")
                    or thread_id
                ),
                "created_at": result_blob.get("timestamp") or session.get("timestamp") or session.get("created"),
                "result": result_blob,
            }
        elif analyses:
            selected = analyses[0]

    if not isinstance(selected, dict):
        return None, (jsonify({"error": "No scorecard is available for this thread."}), 404)

    result = selected.get("result") if isinstance(selected.get("result"), dict) else {}
    component_scores = result.get("component_scores") if isinstance(result.get("component_scores"), dict) else {}
    if not component_scores and isinstance(result.get("scores"), dict):
        component_scores = result.get("scores")
    component_scores = component_scores if isinstance(component_scores, dict) else {}

    overrides = result.get("display_overrides") if isinstance(result.get("display_overrides"), dict) else {}
    financial_impact = result.get("financial_impact") if isinstance(result.get("financial_impact"), dict) else {}
    risks = result.get("top_risks")
    if not isinstance(risks, list):
        risks = result.get("risks")
    risks = risks if isinstance(risks, list) else []
    if "top_risks" in overrides and isinstance(overrides.get("top_risks"), list):
        risks = overrides.get("top_risks")
    recommendations = result.get("recommendations") if isinstance(result.get("recommendations"), list) else []
    executive_summary = (
        overrides.get("executive_summary")
        if "executive_summary" in overrides
        else result.get("executive_summary")
    )
    recommended_scenario = (
        overrides.get("recommended_scenario")
        if "recommended_scenario" in overrides
        else result.get("recommended_scenario")
    )
    custom_blocks = overrides.get("custom_blocks") if isinstance(overrides.get("custom_blocks"), list) else []

    payload = {
        "analysis_id": str(selected.get("analysis_id") or thread_id),
        "project_name": _safe_text(
            overrides.get("title") or result.get("project_name") or result.get("name") or session.get("name") or f"Thread {thread_id}",
            255,
        ),
        "jaspen_score": result.get("jaspen_score") or result.get("overall_score") or result.get("score"),
        "score_category": _safe_text(result.get("score_category"), 64) or None,
        "component_scores": component_scores,
        "financial_impact": financial_impact,
        "risks": risks,
        "recommendations": recommendations,
        "updated_at": selected.get("created_at") or result.get("timestamp") or session.get("timestamp"),
        "scenario_variants": _scorecard_variants_for_export(
            session,
            thread_id,
            selected_scorecard_id=scorecard_id,
            user_id=user_id,
        ),
        # Full content for rich exports (Excel/Word): the per-criterion grid, the
        # written summary, and the rubric (for ordering/labels).
        "dimensions": result.get("dimensions") if isinstance(result.get("dimensions"), dict) else {},
        "executive_summary": _safe_text(executive_summary, 4000) or None,
        "rubric": result.get("rubric") if isinstance(result.get("rubric"), dict) else None,
        "top_risks": risks,
        "recommended_scenario": _safe_text(recommended_scenario, 4000) or None,
        "custom_blocks": custom_blocks,
        "display_overrides": overrides,
        # #4 custom colors: carry the user's brand accent into exports.
        "accent_color": (
            str(overrides.get("accent_color") or result.get("_accent_color") or "").strip()
            or None
        ),
    }
    return payload, None


def _accent_hex(scorecard, default="#A0036C"):
    """Validated brand accent for exports, falling back to Jaspen magenta."""
    raw = str((scorecard or {}).get("accent_color") or "").strip()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", raw):
        return raw.upper()
    return default


def _display_value(value):
    if value is None or value == "":
        return "N/A"
    if isinstance(value, float):
        return f"{value:,.2f}"
    return str(value)


def _safe_float(value):
    try:
        numeric = float(value)
        if numeric == numeric:
            return numeric
    except Exception:
        return None
    return None


def _scorecard_component_rows(scorecard):
    """Return the single authoritative dimension set in rubric order.

    Modern scorecards store rubric-aligned detail in ``dimensions`` while some
    records retain legacy ``component_scores`` keys. Mixing both creates
    duplicate slides and misleading zeroes, so dimensions win whenever present.
    """
    component_scores = scorecard.get("component_scores") if isinstance(scorecard.get("component_scores"), dict) else {}
    dimensions = scorecard.get("dimensions") if isinstance(scorecard.get("dimensions"), dict) else {}
    criteria = ((scorecard.get("rubric") or {}).get("criteria") or []) if isinstance(scorecard.get("rubric"), dict) else []
    criterion_by_key = {
        str(item.get("key")): item
        for item in criteria
        if isinstance(item, dict) and item.get("key")
    }

    source_keys = list(dimensions) if dimensions else list(component_scores)
    ordered_keys = [
        str(item.get("key"))
        for item in criteria
        if isinstance(item, dict) and str(item.get("key") or "") in source_keys
    ]
    ordered_keys.extend(str(key) for key in source_keys if str(key) not in ordered_keys)

    rows = []
    for key in ordered_keys:
        dim = dimensions.get(key) if isinstance(dimensions.get(key), dict) else {}
        criterion = criterion_by_key.get(key) or {}
        value = dim.get("score") if dimensions and dim.get("score") is not None else component_scores.get(key)
        rows.append(
            {
                "key": key,
                "label": dim.get("label") or criterion.get("label") or _format_label(key),
                "value": value,
                "rationale": dim.get("rationale") or dim.get("reasoning") or "",
                "is_risk": bool(dim.get("is_risk", criterion.get("is_risk", False))),
            }
        )
    return rows


def _meaningful_financial_items(scorecard):
    financial = scorecard.get("financial_impact") if isinstance(scorecard.get("financial_impact"), dict) else {}
    items = []
    for key, value in financial.items():
        if str(key).lower() in {"numeric", "_numeric"} or isinstance(value, (dict, list, tuple)):
            continue
        if value is None or str(value).strip().lower() in {"", "n/a", "none", "null"}:
            continue
        items.append(f"{_format_label(key)}: {_display_value(value)}")
    return items


def _grid_int(source, key, default):
    """Read one react-grid-layout coordinate, tolerating strings and junk."""
    try:
        value = int(source.get(key))
    except (TypeError, ValueError, AttributeError):
        return default
    return value


def _scorecard_custom_blocks(scorecard):
    raw = scorecard.get("custom_blocks") if isinstance(scorecard.get("custom_blocks"), list) else []
    blocks = []
    for index, item in enumerate(raw[:12]):
        if not isinstance(item, dict):
            continue
        heading = _safe_text(item.get("heading") or item.get("title") or f"Additional context {index + 1}", 160)
        body = _safe_text(item.get("body") or item.get("text"), 5000)
        if heading or body:
            blocks.append({
                "heading": heading or "Additional context",
                "body": body,
                "type": item.get("type") or "text",
                # Canvas position, so exports can place blocks where the user put
                # them rather than appending them all at the end.
                "x": _grid_int(item, "x", 0),
                "y": _grid_int(item, "y", 900 + index),
                "w": max(1, min(12, _grid_int(item, "w", 6))),
            })
    return blocks


# Fallback arrangement used when a card has no saved section_layout. Matches the
# long-standing export look (score beside the summary) so existing scorecards
# keep exporting the way they always have.
_DEFAULT_EXPORT_SECTION_LAYOUT = [
    {"key": "score", "x": 0, "y": 0, "w": 3},
    {"key": "executive", "x": 3, "y": 0, "w": 9},
    {"key": "dimensions", "x": 0, "y": 1, "w": 12},
    {"key": "risks", "x": 0, "y": 2, "w": 6},
    {"key": "scenario", "x": 6, "y": 2, "w": 6},
]


def _scorecard_section_layout(scorecard):
    """The user's canvas arrangement for the built-in sections, keyed by section.

    Falls back to the historical export arrangement when the card predates
    layout persistence (or the user never rearranged anything).
    """
    raw = scorecard.get("display_overrides") if isinstance(scorecard.get("display_overrides"), dict) else {}
    saved = raw.get("section_layout") if isinstance(raw.get("section_layout"), list) else []
    rows = [row for row in saved if isinstance(row, dict) and row.get("key")]
    is_fallback = not rows
    if is_fallback:
        rows = _DEFAULT_EXPORT_SECTION_LAYOUT
    # Marks a layout the user never saved. Custom blocks carry real canvas row
    # numbers, which don't share a coordinate space with the fallback above, so
    # callers append blocks at the end rather than trying to interleave them.
    layout = {"_fallback": is_fallback}
    for index, row in enumerate(rows):
        layout[str(row.get("key"))] = {
            "x": _grid_int(row, "x", 0),
            "y": _grid_int(row, "y", index),
            "w": max(1, min(12, _grid_int(row, "w", 12))),
            "collapsed": bool(row.get("collapsed")),
            "dim_cols": max(1, min(2, _grid_int(row, "dimCols", 2))),
        }
    return layout


def _pack_layout_rows(placed):
    """Group positioned cards into visual rows the way the canvas grid reads.

    `placed` is a list of {y, x, w, flowable}. Cards are sorted top-to-bottom
    then left-to-right and packed into rows of at most 12 grid columns; a card
    that starts below the current row's band opens a new row. This reproduces
    side-by-side vs stacked arrangement without trying to replicate pixel
    heights — reportlab still flows each card's content naturally.
    """
    ordered = sorted(placed, key=lambda item: (item["y"], item["x"]))
    rows = []
    current = []
    current_width = 0
    current_y = None
    for item in ordered:
        starts_new_band = current_y is not None and item["y"] > current_y
        if current and (current_width + item["w"] > 12 or starts_new_band):
            rows.append(current)
            current = []
            current_width = 0
            current_y = None
        current.append(item)
        current_width += item["w"]
        current_y = item["y"] if current_y is None else current_y
    if current:
        rows.append(current)
    return rows


def _scorecard_variants_for_export(session, thread_id, selected_scorecard_id=None, user_id=None):
    if not isinstance(session, dict):
        return []

    if user_id is not None:
        from app.scenarios_store import load_scenarios_data
        from app.scorecards import collect_peer_scorecards
        snapshots = collect_peer_scorecards(
            user_id,
            thread_id,
            legacy_session=session,
            legacy_thread_data=(load_scenarios_data(user_id) or {}).get(thread_id) or {},
        )
        snapshot_state = {'selected_id': selected_scorecard_id}
    else:
        result_blob = session.get("result") if isinstance(session.get("result"), dict) else {}
        if not result_blob:
            return []
        try:
            snapshot_state = _scorecard_snapshot_state(result_blob, thread_id)
        except Exception:
            return []
        snapshots = snapshot_state.get("snapshots") if isinstance(snapshot_state.get("snapshots"), list) else []
    if len(snapshots) <= 1:
        return []

    selected_id = str(selected_scorecard_id or snapshot_state.get("selected_id") or "").strip()
    variants = []
    for snap in snapshots:
        if not isinstance(snap, dict):
            continue
        snap_id = str(snap.get("id") or snap.get("analysis_id") or "").strip()
        if not snap_id:
            continue
        score = _safe_float(snap.get("jaspen_score") or snap.get("overall_score") or snap.get("score"))
        variants.append(
            {
                "id": snap_id,
                "label": _safe_text(snap.get("project_name") or snap.get("name") or snap.get("label") or "Project", 120),
                "is_baseline": False,
                "is_selected": bool(selected_id and snap_id == selected_id),
                "jaspen_score": score,
                "score_category": _safe_text(snap.get("score_category"), 64) or None,
            }
        )

    return variants


def _list_text_items(items, *, fallback):
    if not isinstance(items, list) or not items:
        return [fallback]
    output = []
    for item in items:
        if isinstance(item, dict):
            text = (
                item.get("title")
                or item.get("recommendation")
                or item.get("risk")
                or item.get("description")
            )
        else:
            text = item
        text = str(text or "").strip()
        if text:
            output.append(text)
    return output or [fallback]


def _scorecard_markdown(scorecard, *, org=None):
    project_name = scorecard.get("project_name") or "Untitled Idea"
    component_scores = scorecard.get("component_scores") or {}
    financial_impact = scorecard.get("financial_impact") or {}
    risks = _list_text_items(scorecard.get("risks"), fallback="No key risks recorded.")
    recommendations = _list_text_items(scorecard.get("recommendations"), fallback="No recommendations recorded.")
    generated_at = scorecard.get("updated_at") or _iso_now()

    lines = [
        f"# {project_name}",
        "",
        f"- **Generated**: {generated_at}",
        f"- **Workspace**: {(org.name if org else 'Personal Workspace')}",
        f"- **Jaspen Score**: {_display_value(scorecard.get('jaspen_score'))}",
        f"- **Score Category**: {_display_value(scorecard.get('score_category'))}",
        "",
        "## Component Scores",
    ]
    if component_scores:
        for key, value in component_scores.items():
            lines.append(f"- **{_format_label(key)}**: {_display_value(value)}")
    else:
        lines.append("- No component scores recorded.")

    lines.extend(["", "## Financial Impact"])
    if financial_impact:
        for key, value in financial_impact.items():
            lines.append(f"- **{_format_label(key)}**: {_display_value(value)}")
    else:
        lines.append("- No financial impact data recorded.")

    lines.extend(["", "## Key Risks"])
    lines.extend(f"- {item}" for item in risks)
    lines.extend(["", "## Recommendations + Next Steps"])
    lines.extend(f"- {item}" for item in recommendations)
    return "\n".join(lines)


def _safe_filename_base(value):
    slug = re.sub(r"[^a-z0-9._-]+", "-", str(value or "").strip().lower()).strip("-")
    return slug[:96] or "jaspen-export"


def _scorecard_pdf_bytes(scorecard, *, org=None):
    project_name = scorecard.get("project_name") or "Untitled Idea"
    markdown_fallback = _scorecard_markdown(scorecard, org=org)

    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import landscape, letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Flowable, KeepTogether, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from xml.sax.saxutils import escape
    except Exception:
        return _markdown_to_pdf_bytes(project_name, markdown_fallback)

    try:
        navy = colors.HexColor("#0F172A")
        accent_hex = _accent_hex(scorecard)
        accent = colors.HexColor(accent_hex)
        ink = colors.HexColor("#334155")
        slate = colors.HexColor("#64748B")
        border = colors.HexColor("#E6EAF2")
        track = colors.HexColor("#EEF2F6")
        page_bg = colors.HexColor("#F8FAFC")

        component_rows = _scorecard_component_rows(scorecard)
        risks = _list_text_items(scorecard.get("risks"), fallback="")
        score_text = _display_value(scorecard.get("jaspen_score"))
        category_text = str(scorecard.get("score_category") or "Not categorized").upper()
        recommended = _safe_text(scorecard.get("recommended_scenario"), 3000)
        if not recommended:
            recommendation_items = _list_text_items(scorecard.get("recommendations"), fallback="")
            recommended = recommendation_items[0] if recommendation_items and recommendation_items[0] else ""
        custom_blocks = _scorecard_custom_blocks(scorecard)

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "WorkspaceTitle",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=22,
            leading=27,
            textColor=navy,
        )
        section_style = ParagraphStyle(
            "WorkspaceSection",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=8.5,
            leading=10,
            textColor=slate,
            uppercase=True,
            tracking=0.7,
        )
        body_style = ParagraphStyle(
            "WorkspaceBody",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=10.5,
            leading=15,
            textColor=ink,
        )
        bullet_style = ParagraphStyle(
            "WorkspaceBullet",
            parent=body_style,
            leftIndent=9,
            firstLineIndent=-7,
            spaceAfter=4,
        )
        small_style = ParagraphStyle(
            "WorkspaceSmall",
            parent=body_style,
            fontSize=8.5,
            leading=11,
            textColor=slate,
        )

        class ScoreBar(Flowable):
            def __init__(self, value, width, color):
                super().__init__()
                self.width = width
                self.height = 5
                numeric = _safe_float(value)
                normalized = (numeric or 0) / 10 if (numeric or 0) > 10 else (numeric or 0)
                self.percent = max(0, min(1, normalized / 10))
                self.color = color

            def draw(self):
                self.canv.setFillColor(track)
                self.canv.roundRect(0, 0, self.width, self.height, 2.5, fill=1, stroke=0)
                if self.percent:
                    self.canv.setFillColor(self.color)
                    self.canv.roundRect(0, 0, self.width * self.percent, self.height, 2.5, fill=1, stroke=0)

        def section_card(title, content, *, left_accent=False, background=colors.white):
            inner = [[Paragraph(escape(str(title).upper()), section_style)], [content]]
            table = Table(inner, colWidths=[None], hAlign="LEFT")
            commands = [
                ("BACKGROUND", (0, 0), (-1, -1), background),
                ("BOX", (0, 0), (-1, -1), 0.75, border),
                ("LINEBELOW", (0, 0), (-1, 0), 0.6, border),
                ("LEFTPADDING", (0, 0), (-1, -1), 12),
                ("RIGHTPADDING", (0, 0), (-1, -1), 12),
                ("TOPPADDING", (0, 0), (-1, 0), 9),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("TOPPADDING", (0, 1), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 1), (-1, -1), 12),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
            if left_accent:
                commands.append(("LINEBEFORE", (0, 0), (0, -1), 2.5, accent))
            table.setStyle(TableStyle(commands))
            return table

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=landscape(letter),
            leftMargin=0.5 * inch,
            rightMargin=0.5 * inch,
            topMargin=0.42 * inch,
            bottomMargin=0.42 * inch,
            title=_safe_text(project_name, 200),
        )

        story = [Paragraph(escape(_safe_text(project_name, 240)), title_style), Spacer(1, 14)]

        score_style = ParagraphStyle(
            "WorkspaceScore", parent=styles["Normal"], fontName="Helvetica-Bold",
            fontSize=34, leading=38, alignment=1, textColor=navy,
        )
        score_meta_style = ParagraphStyle(
            "WorkspaceScoreMeta", parent=small_style, alignment=1,
            fontName="Helvetica-Bold", textColor=accent,
        )
        score_content = [
            Paragraph(escape(score_text), score_style),
            Paragraph(escape(category_text), score_meta_style),
            Spacer(1, 5),
            Paragraph("Strategy scorecard", ParagraphStyle("WorkspaceScoreLabel", parent=body_style, alignment=1, fontName="Helvetica-Bold")),
            Paragraph("Scores reflect Jaspen's analysis.", ParagraphStyle("WorkspaceScoreHint", parent=small_style, alignment=1)),
        ]
        executive = Paragraph(
            escape(_safe_text(scorecard.get("executive_summary") or "No executive summary recorded.", 3000)).replace("\n", "<br/>"),
            body_style,
        )

        # The canvas arrangement the user dragged into place. Everything below is
        # positioned from this rather than a fixed order, so the PDF matches the
        # page instead of guessing.
        layout = _scorecard_section_layout(scorecard)
        dim_cols = (layout.get("dimensions") or {}).get("dim_cols", 2)

        # Dimension columns have to fit INSIDE the Dimensions card, which is
        # itself only as wide as the user made it. Deriving the width instead of
        # hard-coding it keeps the inner grid from overflowing its cell — an
        # overflow makes reportlab mis-measure the row and bump later cards onto
        # a mostly-empty next page.
        content_width = 9.7 * inch
        dim_gutter = 14
        dim_card_width = content_width * (((layout.get("dimensions") or {}).get("w", 12)) / 12) - 24
        dim_col_width = max(1.5 * inch, (dim_card_width - (dim_gutter if dim_cols == 2 else 0)) / dim_cols)
        score_label_width = 0.55 * inch

        dimension_cells = [[]] if dim_cols == 1 else [[], []]
        for index, item in enumerate(component_rows):
            numeric = _safe_float(item.get("value"))
            normalized = numeric / 10 if numeric is not None and numeric > 10 else numeric
            score_label = f"{normalized:.1f}/10" if normalized is not None else "N/A"
            heading = Table(
                [[Paragraph(escape(_safe_text(item.get("label"), 180)), body_style), Paragraph(score_label, small_style)]],
                colWidths=[dim_col_width - score_label_width, score_label_width],
            )
            heading.setStyle(TableStyle([("ALIGN", (1, 0), (1, 0), "RIGHT"), ("VALIGN", (0, 0), (-1, -1), "BOTTOM"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0), ("TOPPADDING", (0, 0), (-1, -1), 0), ("BOTTOMPADDING", (0, 0), (-1, -1), 4)]))
            bar_color = colors.HexColor("#F59E0B") if item.get("is_risk") else navy
            cell = [heading, ScoreBar(item.get("value"), dim_col_width, bar_color), Spacer(1, 10)]
            dimension_cells[index % len(dimension_cells)].extend(cell)
        if not component_rows:
            dimension_cells[0].append(Paragraph("No dimension scores recorded.", body_style))
        dimensions_grid = Table([dimension_cells], colWidths=[dim_col_width] * len(dimension_cells), hAlign="LEFT")
        dimensions_grid.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (0, -1), dim_gutter), ("RIGHTPADDING", (-1, 0), (-1, -1), 0), ("TOPPADDING", (0, 0), (-1, -1), 0), ("BOTTOMPADDING", (0, 0), (-1, -1), 0)]))

        # Every card the canvas shows, tagged with where the user put it.
        def _place(key, flowable, fallback_y):
            spec = layout.get(key)
            # A saved layout lists every section the canvas shows. If the user's
            # layout doesn't mention this one it isn't on their canvas, so leave
            # it out rather than guessing a position for it.
            if spec is None and not layout.get("_fallback"):
                return
            if spec and spec.get("collapsed"):
                return
            placed.append({
                "y": spec.get("y", fallback_y) if spec else fallback_y,
                "x": spec.get("x", 0) if spec else 0,
                "w": spec.get("w", 12) if spec else 12,
                "flowable": flowable,
            })

        placed = []
        _place("score", section_card("Score", score_content), 0)
        _place("executive", section_card("Executive summary", executive), 0)
        _place("dimensions", section_card("Dimensions", dimensions_grid), 1)
        if any(risks):
            risk_content = [Paragraph(f"• {escape(_safe_text(item, 700))}", bullet_style) for item in risks[:8] if item]
            _place("risks", section_card("Top risks", risk_content), 2)
        if recommended:
            _place(
                "scenario",
                section_card("Recommended scenario", Paragraph(escape(recommended).replace("\n", "<br/>"), body_style), left_accent=True),
                2,
            )
        for block_index, block in enumerate(custom_blocks):
            placed.append({
                # With no saved layout there is no shared coordinate space, so
                # keep the historical behaviour: blocks follow the built-ins.
                "y": 900 + block_index if layout.get("_fallback") else block.get("y", 900),
                "x": block.get("x", 0),
                "w": block.get("w", 6),
                "flowable": section_card(
                    block["heading"],
                    Paragraph(escape(block["body"]).replace("\n", "<br/>"), body_style),
                    left_accent=block.get("type") == "callout",
                ),
            })

        gutter = 10
        for row in _pack_layout_rows(placed):
            # Always divide by the full 12-column grid, never by the row's own
            # total: a lone half-width card should stay half width, exactly as
            # it sits on the canvas, rather than stretching across the page.
            widths = [content_width * (item["w"] / 12) - (gutter if index < len(row) - 1 else 0)
                      for index, item in enumerate(row)]
            grid = Table([[item["flowable"] for item in row]], colWidths=widths, hAlign="LEFT")
            # The gutter is already subtracted from each width, so zero the cell
            # padding rather than adding it twice.
            grid.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0), ("TOPPADDING", (0, 0), (-1, -1), 0), ("BOTTOMPADDING", (0, 0), (-1, -1), 0)]))
            # KeepTogether only where it earns its keep — holding side-by-side
            # cards on one page. Wrapping a lone full-width card (a long
            # Dimensions grid) risks a block taller than the page.
            story.extend([KeepTogether(grid) if len(row) > 1 else grid, Spacer(1, 12)])

        def paint_page(canvas, _doc):
            canvas.saveState()
            canvas.setFillColor(page_bg)
            canvas.rect(0, 0, landscape(letter)[0], landscape(letter)[1], fill=1, stroke=0)
            canvas.restoreState()

        doc.build(story, onFirstPage=paint_page, onLaterPages=paint_page)
        buffer.seek(0)
        return buffer.read()
    except Exception:
        return _markdown_to_pdf_bytes(project_name, markdown_fallback)


def _display_chat_role(role):
    normalized = str(role or "").strip().lower()
    if normalized in {"assistant", "ai", "bot"}:
        return "Jaspen"
    return "You"


def _transcript_markdown(session, *, org=None):
    if not isinstance(session, dict):
        return ""
    thread_name = _safe_text(session.get("name") or "Conversation Transcript", 255)
    workspace_name = org.name if org else "Personal Workspace"
    generated_at = session.get("timestamp") or session.get("created") or _iso_now()
    chat_history = session.get("chat_history") if isinstance(session.get("chat_history"), list) else []

    lines = [
        f"# {thread_name}",
        "",
        f"- **Generated**: {generated_at}",
        f"- **Workspace**: {workspace_name}",
        f"- **Thread ID**: {str(session.get('session_id') or '').strip() or 'N/A'}",
        "",
        "## Conversation",
    ]

    if not chat_history:
        lines.append("")
        lines.append("_No conversation messages available._")
        return "\n".join(lines)

    for message in chat_history:
        if not isinstance(message, dict):
            continue
        content = _safe_text(message.get("content") or message.get("text") or "", 20000)
        attachments = message.get("attachments") if isinstance(message.get("attachments"), list) else []
        if not str(content).strip() and not attachments:
            continue
        role_label = _display_chat_role(message.get("role"))
        timestamp = str(message.get("timestamp") or "").strip()
        lines.extend(["", f"### {role_label}"])
        if timestamp:
            lines.append(f"_Timestamp: {timestamp}_")
            lines.append("")
        if str(content).strip():
            lines.append(str(content).strip())
        if attachments:
            lines.append("")
            lines.append("Attached files:")
            for attachment in attachments:
                if not isinstance(attachment, dict):
                    continue
                name = _safe_text(attachment.get("name") or "attachment", 255)
                attachment_type = _safe_text(attachment.get("type") or attachment.get("kind") or "file", 120)
                size = int(attachment.get("size") or 0)
                size_kb = max(1, round(size / 1024)) if size > 0 else None
                suffix = f" ({attachment_type}{f', {size_kb} KB' if size_kb else ''})"
                lines.append(f"- {name}{suffix}")

    return "\n".join(lines).strip()


def _send_bytes(payload, *, filename, mimetype):
    buffer = io.BytesIO(payload)
    buffer.seek(0)
    return send_file(buffer, mimetype=mimetype, as_attachment=True, download_name=filename)


def _pptx_bytes(scorecard, *, org=None, peers=None):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError("python-pptx is required for PowerPoint export.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(0x16, 0x1F, 0x3B)
    accent_hex = _accent_hex(scorecard)
    magenta = RGBColor.from_string(accent_hex.lstrip("#"))
    gray = RGBColor(0x60, 0x67, 0x74)
    pale = RGBColor(0xEF, 0xF9, 0xFC)
    border = RGBColor(0xD7, 0xDE, 0xE8)

    def add_text(
        slide,
        text,
        left,
        top,
        width,
        height,
        *,
        size=18,
        color=navy,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.01)
        frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = valign
        para = frame.paragraphs[0]
        para.text = _safe_text(text, 1600)
        para.alignment = align
        para.font.name = "Aptos"
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = color
        return box

    def add_rule(slide, left, top, width, *, color=border, height=0.02):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_header(slide, title, subtitle=None):
        add_rule(slide, 0.65, 0.38, 0.78, color=magenta, height=0.06)
        add_text(slide, title, 0.65, 0.55, 12.0, 0.62, size=27, bold=True)
        if subtitle:
            add_text(slide, subtitle, 0.65, 1.17, 12.0, 0.35, size=11, color=gray)

    def add_list(slide, items, left, top, width, height, *, size=15, fallback):
        cleaned = _list_text_items(items, fallback=fallback)[:6]
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        for idx, item in enumerate(cleaned):
            para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            para.text = f"• {_safe_text(item, 280)}"
            para.font.name = "Aptos"
            para.font.size = Pt(size)
            para.font.color.rgb = navy
            para.space_after = Pt(8)
        return box

    def component_display(value):
        numeric = _safe_float(value)
        if numeric is None:
            return "N/A", 0
        normalized = numeric / 10 if numeric > 10 else numeric
        return f"{normalized:.1f}/10", max(0, min(100, normalized * 10))

    project_name = scorecard.get("project_name") or "Untitled Idea"
    generated_at = scorecard.get("updated_at") or _iso_now()
    workspace_name = org.name if org else "Personal Workspace"
    category = scorecard.get("score_category") or "Not categorized"
    component_rows = [
        (row.get("label"), row.get("value"), row.get("rationale"))
        for row in _scorecard_component_rows(scorecard)
    ]
    if not component_rows:
        component_rows = [("No component scores recorded", None, "")]

    # Slide 1: the decision at a glance.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, project_name, f"{workspace_name}  •  Generated {generated_at}")
    score_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.65), Inches(1.75), Inches(3.0), Inches(4.85),
    )
    score_panel.fill.solid()
    score_panel.fill.fore_color.rgb = pale
    score_panel.line.color.rgb = border
    add_text(slide, "JASPEN SCORE", 0.95, 2.1, 2.4, 0.35, size=12, color=gray, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, _display_value(scorecard.get("jaspen_score")), 0.9, 2.55, 2.5, 1.05, size=48, color=magenta, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, str(category).upper(), 0.95, 3.65, 2.4, 0.35, size=13, color=magenta, bold=True, align=PP_ALIGN.CENTER)
    add_rule(slide, 1.15, 4.25, 2.0)
    add_text(slide, "Editable scorecard", 0.95, 4.6, 2.4, 0.32, size=12, color=gray, align=PP_ALIGN.CENTER)
    add_text(slide, "Every label, score, bar, and narrative can be updated directly in PowerPoint.", 0.95, 5.05, 2.4, 0.95, size=13, color=navy, align=PP_ALIGN.CENTER)

    add_text(slide, "Executive summary", 4.25, 1.85, 8.25, 0.45, size=18, bold=True)
    add_rule(slide, 4.25, 2.38, 8.25)
    executive_summary = scorecard.get("executive_summary") or "No executive summary recorded."
    add_text(slide, executive_summary, 4.25, 2.65, 8.25, 2.0, size=18)
    add_text(slide, "Decision signal", 4.25, 5.05, 8.25, 0.38, size=14, color=gray, bold=True)
    decision_signal = scorecard.get("recommended_scenario") or (scorecard.get("recommendations") or [None])[0]
    if isinstance(decision_signal, dict):
        decision_signal = decision_signal.get("text") or decision_signal.get("action")
    add_text(slide, decision_signal or "Review the score, evidence, and risks before committing resources.", 4.25, 5.48, 8.25, 1.0, size=18, color=magenta, bold=True)

    # Decision Confidence slides.
    #
    # Deck-shaped, deliberately NOT the emailed report with slide breaks. The
    # email is scrolled by one reader; a deck is read from across a room while
    # someone talks over it. So the deck carries the split, the finding, the
    # single action, and only the criteria that could move the score or the
    # ranking. Everything else lives in the email and the workspace.
    #
    # Failure here must never cost the export: a scorecard without a report
    # simply skips these slides.
    try:
        from app.decision_report import build_report

        # Peers enable the standing line, which no single card can state.
        report = build_report(scorecard, peers=peers)
    except Exception:
        report = None

    if report:
        summary = report.get("summary") or {}
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header(slide, "Decision Confidence", "How much of this decision rests on evidence")

        backed = report.get("evidence_backed_pct") or 0
        assumed = report.get("assumption_dependent_pct") or 0

        add_text(slide, f"{backed}%", 0.65, 1.85, 2.6, 1.0, size=54, color=navy, bold=True)
        add_text(slide, "EVIDENCE-BACKED", 0.65, 2.95, 2.6, 0.3, size=11, color=gray, bold=True)
        add_text(slide, f"{assumed}%", 3.55, 1.85, 2.6, 1.0, size=54, color=magenta, bold=True)
        add_text(slide, "ASSUMPTION-DEPENDENT", 3.55, 2.95, 2.6, 0.3, size=11, color=gray, bold=True)

        # One bar, two shares. Widths are the split, so the figure and the
        # picture cannot disagree.
        bar_w = 5.5
        backed_w = max(0.05, bar_w * (backed / 100.0))
        add_rule(slide, 0.65, 3.45, backed_w, color=RGBColor(0x0E, 0x6B, 0x3F), height=0.12)
        add_rule(slide, 0.65 + backed_w, 3.45, bar_w - backed_w,
                 color=RGBColor(0xF5, 0x9E, 0x0B), height=0.12)
        add_text(slide, "Evidence-backed share of the weighted decision",
                 0.65, 3.65, 5.5, 0.3, size=10, color=gray)

        briefing = " ".join(
            summary[key] for key in ("verdict", "standing", "confidence", "concentration")
            if summary.get(key)
        )
        add_text(slide, briefing, 6.6, 1.85, 6.1, 2.2, size=15)
        if summary.get("sensitivity"):
            add_rule(slide, 6.6, 4.15, 6.1)
            add_text(slide, summary["sensitivity"], 6.6, 4.35, 6.1, 1.1, size=15, bold=True)
        if summary.get("next_step"):
            add_text(slide, "DO THIS NEXT", 0.65, 4.35, 5.5, 0.3, size=11, color=magenta, bold=True)
            add_text(slide, summary["next_step"], 0.65, 4.7, 5.5, 1.6, size=14, color=navy)

        # Only what could move the answer earns a slide. A deck listing every
        # criterion is the email, and nobody reads the email from a projector.
        material = report.get("material") or []
        if material:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_header(slide, "What could change the answer",
                       "Assumptions ranked by their power to move the score")
            top = 1.9
            for criterion in material[:4]:
                add_text(slide, criterion["label"], 0.65, top, 6.4, 0.4, size=17, bold=True)
                add_text(slide,
                         f'{criterion["weight_pct"]}% of the decision  •  '
                         f'{criterion["grade_label"]}  •  {criterion["swing"]} points of exposure',
                         0.65, top + 0.42, 6.4, 0.32, size=12, color=gray)
                if criterion.get("evidence_needed"):
                    add_text(slide, f'Evidence needed: {criterion["evidence_needed"]}',
                             7.3, top, 5.4, 0.9, size=13, color=navy)
                add_rule(slide, 0.65, top + 1.0, 12.0)
                top += 1.25

        # Risk, at slide density rather than register density.
        #
        # Deliberately NOT every field. A slide showing likelihood, impact,
        # category, mitigation, mitigation cost and residual for each of five
        # risks is a spreadsheet photographed badly. The deck carries what a
        # room needs to react to: the risk, what survives mitigation, and the
        # exposure. Mitigation detail and costs live in the email.
        risks = report.get("risks") or []
        if risks:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_header(slide, "Risk register", "Ordered by unmitigated exposure")
            add_text(slide, "RISK", 0.65, 1.75, 8.0, 0.3, size=10, color=gray, bold=True)
            add_text(slide, "IF MITIGATED", 8.9, 1.75, 2.0, 0.3, size=10, color=gray, bold=True)
            add_text(slide, "EXPOSURE", 11.1, 1.75, 1.6, 0.3, size=10, color=gray, bold=True)
            add_rule(slide, 0.65, 2.05, 12.0)

            top = 2.2
            for risk in risks[:5]:
                add_text(slide, risk.get("risk") or "Untitled risk", 0.65, top, 8.0, 0.85, size=14)
                residual = risk.get("residual")
                residual_color = {
                    "High": RGBColor(0x9F, 0x1F, 0x16),
                    "Medium": RGBColor(0x8A, 0x54, 0x06),
                    "Low": RGBColor(0x0E, 0x6B, 0x3F),
                }.get(residual, gray)
                add_text(slide, residual or "Not rated", 8.9, top, 2.0, 0.4,
                         size=14, color=residual_color, bold=True)
                add_text(slide, risk.get("impact") or "Not sized", 11.1, top, 1.6, 0.4,
                         size=14, color=navy)
                add_rule(slide, 0.65, top + 0.9, 12.0, color=RGBColor(0xEF, 0xF1, 0xF6))
                top += 1.05

            footnotes = ["Residual assumes the mitigation is carried out; "
                         "Jaspen does not track whether it has been."]
            if len(risks) > 5:
                footnotes.append(f"{len(risks) - 5} further risks are in the full report.")
            add_text(slide, "  ".join(footnotes), 0.65, top + 0.1, 12.0, 0.4,
                     size=11, color=gray)

        # The provenance limit travels with the deck. A slide that showed
        # Jaspen's reasoning without it would read as a source citation to a
        # room that never saw the caveat.
        add_text(slide, report["provenance_note"], 0.65, 6.75, 12.0, 0.5, size=9, color=gray)

    # Slides 2+: criterion detail in readable groups of six.
    for page_index in range(0, len(component_rows), 6):
        page_rows = component_rows[page_index:page_index + 6]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_number = page_index // 6 + 1
        page_count = (len(component_rows) + 5) // 6
        subtitle = "Scores, rubric labels, and supporting rationale"
        if page_count > 1:
            subtitle = f"{subtitle}  •  {page_number} of {page_count}"
        add_header(slide, "Score breakdown", subtitle)
        for row_index, (label, value, rationale) in enumerate(page_rows):
            y = 1.72 + row_index * 0.88
            score_text, percent = component_display(value)
            add_text(slide, label, 0.7, y, 3.55, 0.48, size=14, bold=True)
            add_text(slide, score_text, 11.75, y, 0.85, 0.35, size=14, color=magenta, bold=True, align=PP_ALIGN.RIGHT)
            add_rule(slide, 4.35, y + 0.08, 7.15, color=border, height=0.14)
            if percent > 0:
                add_rule(slide, 4.35, y + 0.08, 7.15 * percent / 100, color=magenta, height=0.14)
            if rationale:
                add_text(slide, rationale, 0.7, y + 0.52, 11.9, 0.24, size=10.5, color=gray)
            add_rule(slide, 0.7, y + 0.78, 11.9, height=0.01)

    # Decision narrative: include only meaningful saved sections. Custom
    # Workspace blocks (for example, Mitigation) are first-class deck content.
    detail_sections = []
    risk_items = _list_text_items(scorecard.get("risks"), fallback="")
    if any(risk_items):
        detail_sections.append({"title": "Top risks", "items": risk_items})

    recommended = _safe_text(scorecard.get("recommended_scenario"), 2400)
    if not recommended:
        rec_items = _list_text_items(scorecard.get("recommendations"), fallback="")
        if any(rec_items):
            detail_sections.append({"title": "Recommendations + next steps", "items": rec_items, "accent": True})
    else:
        detail_sections.append({"title": "Recommended scenario", "body": recommended, "accent": True})

    for block in _scorecard_custom_blocks(scorecard):
        detail_sections.append({
            "title": block.get("heading") or "Additional context",
            "body": block.get("body") or "",
            "accent": block.get("type") == "callout",
        })

    financial_items = _meaningful_financial_items(scorecard)
    if financial_items:
        detail_sections.append({"title": "Financial impact", "items": financial_items})

    if not detail_sections:
        detail_sections.append({"title": "Decision details", "body": "No additional decision details recorded."})

    for section_index in range(0, len(detail_sections), 2):
        pair = detail_sections[section_index:section_index + 2]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = "Decision details" if section_index == 0 else "Additional context"
        subtitle = "Saved scorecard narrative and supporting context"
        add_header(slide, title, subtitle)
        for column, section in enumerate(pair):
            left = 0.7 if column == 0 else 6.95
            add_text(slide, section["title"], left, 1.78, 5.65, 0.52, size=18, bold=True)
            add_rule(slide, left, 2.3, 5.65, color=magenta if section.get("accent") else border)
            if section.get("items"):
                add_list(slide, section["items"], left, 2.58, 5.65, 3.95, fallback="")
            else:
                add_text(slide, section.get("body") or "", left, 2.58, 5.65, 3.95, size=16)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _wbs_csv_bytes(project_wbs):
    tasks = project_wbs.get("tasks") if isinstance(project_wbs, dict) else []
    if not isinstance(tasks, list) or not tasks:
        return None

    output = io.StringIO()
    writer = csv.DictWriter(
        output,
        fieldnames=[
            "id",
            "title",
            "status",
            "owner",
            "suggested_role",
            "phase",
            "priority",
            "due_date",
            "timeline_days",
            "estimated_days",
            "depends_on",
            "jira_issue_key",
            "description",
            "rationale",
            "risk_area",
            "order",
        ],
    )
    writer.writeheader()
    for task in tasks:
        if not isinstance(task, dict):
            continue
        refs = task.get("external_refs") if isinstance(task.get("external_refs"), dict) else {}
        writer.writerow(
            {
                "id": task.get("id"),
                "title": task.get("title"),
                "status": task.get("status"),
                "owner": task.get("owner"),
                "suggested_role": task.get("suggested_role"),
                "phase": task.get("phase"),
                "priority": task.get("priority"),
                "due_date": task.get("due_date"),
                "timeline_days": task.get("timeline_days"),
                "estimated_days": task.get("estimated_days"),
                "depends_on": ",".join(task.get("depends_on") or []),
                "jira_issue_key": refs.get("jira_issue_key"),
                "description": task.get("description"),
                "rationale": task.get("rationale"),
                "risk_area": task.get("risk_area"),
                "order": task.get("order"),
            }
        )
    return output.getvalue().encode("utf-8")


def _xlsx_text(value):
    """Keep exported user text from being interpreted as an Excel formula."""
    text = str(value or "")
    if text.startswith(("=", "+", "-", "@")):
        return f"'{text}"
    return text


def _finalize_xlsx_text_cells(worksheet):
    """Hide the formula-injection escape while retaining a literal text cell."""
    for row in worksheet.iter_rows():
        for cell in row:
            value = cell.value
            if not isinstance(value, str) or len(value) < 2:
                continue
            if value[0] == "'" and value[1] in "=+-@":
                cell.value = value[1:]
                cell.data_type = "s"
                cell.quotePrefix = True


def _xlsx_date(value):
    text = str(value or "").strip()
    if not text:
        return None
    try:
        return datetime.strptime(text[:10], "%Y-%m-%d").date()
    except (TypeError, ValueError):
        return _xlsx_text(text)


def _wbs_phase_names(project_wbs, tasks):
    names = []
    seen = set()
    for phase in project_wbs.get("phases") or []:
        if not isinstance(phase, dict):
            continue
        name = str(phase.get("name") or phase.get("title") or "").strip()
        if name and name not in seen:
            names.append(name)
            seen.add(name)
    for task in tasks:
        if not isinstance(task, dict):
            continue
        name = str(task.get("phase") or "Execution").strip() or "Execution"
        if name not in seen:
            names.append(name)
            seen.add(name)
    return names


def _wbs_phase_date_range(tasks, phase_name):
    phase_tasks = [
        task
        for task in tasks
        if str(task.get("phase") or "Execution").strip() == phase_name
    ]
    start_dates = [_xlsx_date(task.get("start_date")) for task in phase_tasks]
    due_dates = [_xlsx_date(task.get("due_date")) for task in phase_tasks]
    start_dates = [value for value in start_dates if isinstance(value, date)]
    due_dates = [value for value in due_dates if isinstance(value, date)]
    if not start_dates or not due_dates:
        return ""
    return f"{min(start_dates):%Y-%m-%d} to {max(due_dates):%Y-%m-%d}"


def _wbs_xlsx_bytes(project_wbs, *, project_name=None, workspace_name=None):
    tasks = project_wbs.get("tasks") if isinstance(project_wbs, dict) else []
    tasks = [task for task in tasks if isinstance(task, dict) and str(task.get("title") or "").strip()]
    if not tasks:
        return None

    try:
        from openpyxl import Workbook
        from openpyxl.formatting.rule import FormulaRule
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.worksheet.datavalidation import DataValidation
        from openpyxl.worksheet.table import Table, TableStyleInfo
    except Exception as exc:
        raise RuntimeError("openpyxl is required for Excel export.") from exc

    wb = Workbook()
    wb.properties.title = _xlsx_text(project_name or project_wbs.get("name") or "Execution Plan")
    wb.properties.subject = "Jaspen execution plan export"
    wb.properties.creator = "Jaspen"
    try:
        wb.calculation.calcMode = "auto"
        wb.calculation.fullCalcOnLoad = True
        wb.calculation.forceFullCalc = True
    except Exception:
        pass

    overview = wb.active
    overview.title = "Overview"
    task_sheet = wb.create_sheet("Tasks")

    navy = "161F3B"
    rose = "A0036C"
    slate = "475569"
    pale = "F4F6FA"
    line = "D9E1EC"
    white = "FFFFFF"
    green = "DCFCE7"
    blue = "DBEAFE"
    amber = "FEF3C7"
    red = "FEE2E2"
    thin = Side(style="thin", color=line)
    section_border = Border(bottom=Side(style="medium", color=rose))
    body_border = Border(bottom=thin)
    title = _xlsx_text(project_name or project_wbs.get("name") or "Execution Plan")
    workspace = _xlsx_text(workspace_name or "Personal Workspace")
    generated_at = datetime.now(timezone.utc).replace(tzinfo=None)
    task_data_last_row = len(tasks) + 1
    task_phase_range = f"'Tasks'!$B$2:$B${task_data_last_row}"
    task_title_range = f"'Tasks'!$C$2:$C${task_data_last_row}"
    task_status_range = f"'Tasks'!$D$2:$D${task_data_last_row}"
    task_priority_range = f"'Tasks'!$E$2:$E${task_data_last_row}"

    # Overview sheet: a compact project snapshot plus phase-level rollup.
    overview.merge_cells("A1:F1")
    overview["A1"] = title
    overview["A1"].font = Font(bold=True, size=18, color=white)
    overview["A1"].fill = PatternFill("solid", fgColor=navy)
    overview["A1"].alignment = Alignment(vertical="center")
    overview.row_dimensions[1].height = 30

    overview.merge_cells("A2:F2")
    overview["A2"] = "Use the Tasks tab filters to focus on a phase, owner, status, priority, or date range."
    overview["A2"].font = Font(size=10, color=slate)
    overview["A2"].fill = PatternFill("solid", fgColor=pale)
    overview["A2"].alignment = Alignment(wrap_text=True, vertical="center")
    overview.row_dimensions[2].height = 30

    metadata = [
        ("Project", title),
        ("Workspace", workspace),
        ("Plan start", _xlsx_date(project_wbs.get("start_date"))),
        ("Generated", generated_at),
    ]
    for row_idx, (label, value) in enumerate(metadata, start=4):
        overview.cell(row=row_idx, column=1, value=label).font = Font(bold=True, color=navy)
        overview.cell(row=row_idx, column=2, value=value)
        overview.cell(row=row_idx, column=1).border = body_border
        overview.cell(row=row_idx, column=2).border = body_border
    overview["B6"].number_format = "yyyy-mm-dd"
    overview["B7"].number_format = "yyyy-mm-dd hh:mm"

    overview["A9"] = "Progress snapshot"
    overview["A9"].font = Font(bold=True, size=12, color=navy)
    overview["A9"].border = section_border
    overview.merge_cells("A9:B9")
    summary_rows = [
        ("Total tasks", f"=COUNTA({task_title_range})", "0"),
        ("To Do", f'=COUNTIF({task_status_range},"To Do")', "0"),
        ("In Progress", f'=COUNTIF({task_status_range},"In Progress")', "0"),
        ("Blocked", f'=COUNTIF({task_status_range},"Blocked")', "0"),
        ("Done", f'=COUNTIF({task_status_range},"Done")', "0"),
        ("Completion", f'=IFERROR(COUNTIF({task_status_range},"Done")/COUNTA({task_title_range}),0)', "0%"),
    ]
    for row_idx, (label, formula, number_format) in enumerate(summary_rows, start=10):
        overview.cell(row=row_idx, column=1, value=label).font = Font(color=slate)
        overview.cell(row=row_idx, column=2, value=formula).font = Font(bold=True, color=rose)
        overview.cell(row=row_idx, column=2).number_format = number_format

    phases = _wbs_phase_names(project_wbs, tasks)
    phase_start = 18
    overview.cell(row=phase_start, column=1, value="Phase summary")
    overview.cell(row=phase_start, column=1).font = Font(bold=True, size=12, color=navy)
    overview.cell(row=phase_start, column=1).border = section_border
    overview.merge_cells(start_row=phase_start, start_column=1, end_row=phase_start, end_column=6)
    phase_headers = ["Phase #", "Phase", "Tasks", "Done", "High Priority", "Date Range"]
    for col_idx, header in enumerate(phase_headers, start=1):
        cell = overview.cell(row=phase_start + 1, column=col_idx, value=header)
        cell.font = Font(bold=True, color=white)
        cell.fill = PatternFill("solid", fgColor=navy)
        cell.alignment = Alignment(vertical="center")
    for offset, phase_name in enumerate(phases, start=1):
        row_idx = phase_start + 1 + offset
        overview.cell(row=row_idx, column=1, value=offset)
        overview.cell(row=row_idx, column=2, value=_xlsx_text(phase_name))
        overview.cell(row=row_idx, column=3, value=f'=COUNTIF({task_phase_range},B{row_idx})')
        overview.cell(row=row_idx, column=4, value=f'=COUNTIFS({task_phase_range},B{row_idx},{task_status_range},"Done")')
        overview.cell(row=row_idx, column=5, value=f'=COUNTIFS({task_phase_range},B{row_idx},{task_priority_range},"High")')
        overview.cell(row=row_idx, column=6, value=_wbs_phase_date_range(tasks, phase_name))
        for col_idx in range(1, 7):
            overview.cell(row=row_idx, column=col_idx).border = body_border
        overview.cell(row=row_idx, column=6).alignment = Alignment(wrap_text=True)

    overview.sheet_view.showGridLines = False
    for column, width in {"A": 18, "B": 32, "C": 12, "D": 12, "E": 16, "F": 25}.items():
        overview.column_dimensions[column].width = width

    # Tasks sheet: one flat, filterable row per task. Phase rows are data, not
    # decorative headers, so filtering Phase # = 1 returns every Phase 1 task.
    headers = [
        "Phase #",
        "Phase",
        "Task",
        "Status",
        "Priority",
        "Owner",
        "Suggested Role",
        "Start Date",
        "Due Date",
        "Duration (days)",
        "Dependencies",
        "Description / Acceptance Criteria",
        "Risk Area",
        "Rationale",
        "Jira Key",
        "Function",
        "Activity Type",
        "Task ID",
        "Dependency IDs",
        "Task Order",
    ]
    task_sheet.append(headers)
    phase_order = {name: index for index, name in enumerate(phases, start=1)}
    task_by_id = {str(task.get("id") or ""): task for task in tasks}

    def task_sort_key(task):
        phase_name = str(task.get("phase") or "Execution").strip() or "Execution"
        order = task.get("order")
        try:
            order = int(order)
        except (TypeError, ValueError):
            order = tasks.index(task) + 1
        return (phase_order.get(phase_name, len(phase_order) + 1), order)

    status_labels = {
        "todo": "To Do",
        "in_progress": "In Progress",
        "blocked": "Blocked",
        "done": "Done",
    }
    for fallback_order, task in enumerate(sorted(tasks, key=task_sort_key), start=1):
        phase_name = str(task.get("phase") or "Execution").strip() or "Execution"
        raw_dependency_ids = task.get("depends_on") if isinstance(task.get("depends_on"), list) else []
        dependency_ids = [str(dep).strip() for dep in raw_dependency_ids if str(dep).strip()]
        dependency_names = [
            str(task_by_id.get(dep_id, {}).get("title") or dep_id).strip()
            for dep_id in dependency_ids
        ]
        refs = task.get("external_refs") if isinstance(task.get("external_refs"), dict) else {}
        status_key = str(task.get("status") or "todo").strip().lower()
        priority = str(task.get("priority") or "").strip().lower()
        task_order = task.get("order")
        try:
            task_order = int(task_order)
        except (TypeError, ValueError):
            task_order = fallback_order
        duration = task.get("timeline_days")
        if duration is None:
            duration = task.get("estimated_days")
        try:
            duration = int(duration) if duration is not None else None
        except (TypeError, ValueError):
            duration = None
        task_sheet.append([
            phase_order.get(phase_name, len(phase_order) + 1),
            _xlsx_text(phase_name),
            _xlsx_text(task.get("title")),
            status_labels.get(status_key, _format_label(status_key)),
            _format_label(priority) if priority else "",
            _xlsx_text(task.get("owner")),
            _xlsx_text(task.get("suggested_role")),
            _xlsx_date(task.get("start_date")),
            _xlsx_date(task.get("due_date")),
            duration,
            _xlsx_text("; ".join(dependency_names)),
            _xlsx_text(task.get("description") or task.get("acceptance")),
            _xlsx_text(task.get("risk_area")),
            _xlsx_text(task.get("rationale")),
            _xlsx_text(task.get("jira_issue_key") or refs.get("jira_issue_key")),
            _xlsx_text(task.get("function")),
            _xlsx_text(task.get("activity_type")),
            _xlsx_text(task.get("id")),
            _xlsx_text("; ".join(dependency_ids)),
            task_order,
        ])

    last_row = task_sheet.max_row
    last_col = task_sheet.max_column
    table = Table(displayName="ExecutionTasks", ref=f"A1:T{last_row}")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    task_sheet.add_table(table)
    task_sheet.sheet_view.showGridLines = False
    task_sheet.row_dimensions[1].height = 28
    for cell in task_sheet[1]:
        cell.font = Font(bold=True, color=white)
        cell.fill = PatternFill("solid", fgColor=navy)
        cell.alignment = Alignment(wrap_text=True, vertical="center")

    for row_idx in range(2, last_row + 1):
        for col_idx in range(1, last_col + 1):
            task_sheet.cell(row=row_idx, column=col_idx).alignment = Alignment(
                wrap_text=col_idx in {3, 11, 12, 13, 14, 19},
                vertical="top",
            )
        task_sheet.cell(row=row_idx, column=8).number_format = "yyyy-mm-dd"
        task_sheet.cell(row=row_idx, column=9).number_format = "yyyy-mm-dd"
        task_sheet.cell(row=row_idx, column=10).number_format = "0"

    status_validation = DataValidation(type="list", formula1='"To Do,In Progress,Blocked,Done"', allow_blank=False)
    priority_validation = DataValidation(type="list", formula1='"High,Medium,Low"', allow_blank=True)
    task_sheet.add_data_validation(status_validation)
    task_sheet.add_data_validation(priority_validation)
    status_validation.add(f"D2:D{last_row}")
    priority_validation.add(f"E2:E{last_row}")

    status_range = f"D2:D{last_row}"
    priority_range = f"E2:E{last_row}"
    for label, color in (("Done", green), ("In Progress", blue), ("Blocked", red), ("To Do", pale)):
        task_sheet.conditional_formatting.add(
            status_range,
            FormulaRule(formula=[f'$D2="{label}"'], fill=PatternFill("solid", fgColor=color)),
        )
    for label, color in (("High", red), ("Medium", amber), ("Low", green)):
        task_sheet.conditional_formatting.add(
            priority_range,
            FormulaRule(formula=[f'$E2="{label}"'], fill=PatternFill("solid", fgColor=color)),
        )

    widths = {
        "A": 10, "B": 24, "C": 38, "D": 14, "E": 11, "F": 20, "G": 20,
        "H": 12, "I": 12, "J": 14, "K": 34, "L": 52, "M": 20, "N": 44,
        "O": 14, "P": 20, "Q": 16, "R": 22, "S": 30, "T": 12,
    }
    for column, width in widths.items():
        task_sheet.column_dimensions[column].width = width
    task_sheet.print_title_rows = "1:1"
    task_sheet.page_setup.orientation = "landscape"
    task_sheet.page_setup.fitToWidth = 1
    task_sheet.sheet_properties.pageSetUpPr.fitToPage = True

    _finalize_xlsx_text_cells(overview)
    _finalize_xlsx_text_cells(task_sheet)

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def _scorecard_dim_keys(scorecard):
    dims = scorecard.get("dimensions") if isinstance(scorecard.get("dimensions"), dict) else {}
    rubric = scorecard.get("rubric") if isinstance(scorecard.get("rubric"), dict) else {}
    ordered = [c.get("key") for c in (rubric.get("criteria") or []) if isinstance(c, dict) and c.get("key")]
    keys = [k for k in ordered if k in dims] or list(dims.keys())
    return dims, keys


def _scorecard_risk_texts(scorecard):
    risks = scorecard.get("top_risks") if isinstance(scorecard.get("top_risks"), list) else []
    out = []
    for r in risks:
        t = r if isinstance(r, str) else (r.get("risk") or r.get("text") or "" if isinstance(r, dict) else "")
        if str(t).strip():
            out.append(str(t).strip())
    return out


def _xlsx_bytes(scorecard, *, org=None):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except Exception as exc:
        raise RuntimeError("openpyxl is required for Excel export.") from exc

    wb = Workbook()
    ws = wb.active
    ws.title = "Scorecard"
    navy = "161F3B"
    hdr_font = Font(bold=True, color="FFFFFF", size=11)
    hdr_fill = PatternFill("solid", fgColor=navy)
    thin = Side(style="thin", color="DBE3EE")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    wrap = Alignment(wrap_text=True, vertical="top")

    name = scorecard.get("project_name") or "Untitled Idea"
    ws["A1"] = name
    ws["A1"].font = Font(bold=True, size=16, color=navy)
    ws["A2"] = f"{(org.name if org else 'Personal Workspace')} · Generated {scorecard.get('updated_at') or _iso_now()}"
    ws["A2"].font = Font(size=9, color="6B7280")
    ws["A3"] = "Overall Score"
    ws["A3"].font = Font(bold=True)
    ws["B3"] = scorecard.get("jaspen_score")
    ws["B3"].font = Font(bold=True, color="A0036C")
    if scorecard.get("score_category"):
        ws["C3"] = scorecard.get("score_category")

    dims, keys = _scorecard_dim_keys(scorecard)
    row = 5
    for c, h in enumerate(["Criterion", "Score (0-100)", "Confidence", "Rationale"], start=1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.border = border
    row += 1
    for k in keys:
        d = dims.get(k) if isinstance(dims.get(k), dict) else {}
        ws.cell(row=row, column=1, value=d.get("label") or k).border = border
        ws.cell(row=row, column=2, value=d.get("score")).border = border
        ws.cell(row=row, column=3, value=str(d.get("confidence") or "")).border = border
        rc = ws.cell(row=row, column=4, value=str(d.get("rationale") or ""))
        rc.border = border
        rc.alignment = wrap
        row += 1

    if scorecard.get("executive_summary"):
        row += 1
        ws.cell(row=row, column=1, value="Executive Summary").font = Font(bold=True, color=navy)
        row += 1
        ws.cell(row=row, column=1, value=str(scorecard.get("executive_summary"))).alignment = wrap
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
        row += 2

    risk_texts = _scorecard_risk_texts(scorecard)
    if risk_texts:
        ws.cell(row=row, column=1, value="Top Risks").font = Font(bold=True, color=navy)
        row += 1
        for t in risk_texts:
            ws.cell(row=row, column=1, value=f"• {t}").alignment = wrap
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
            row += 1

    for col, width in {"A": 28, "B": 14, "C": 14, "D": 64}.items():
        ws.column_dimensions[col].width = width

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _docx_bytes(scorecard, *, org=None):
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except Exception as exc:
        raise RuntimeError("python-docx is required for Word export.") from exc

    doc = Document()
    _ah = _accent_hex(scorecard)
    magenta = RGBColor(int(_ah[1:3], 16), int(_ah[3:5], 16), int(_ah[5:7], 16))
    name = scorecard.get("project_name") or "Untitled Idea"
    doc.add_heading(name, level=0)
    sub = doc.add_paragraph(f"{(org.name if org else 'Personal Workspace')} · Generated {scorecard.get('updated_at') or _iso_now()}")
    if sub.runs:
        sub.runs[0].font.size = Pt(9)
        sub.runs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    p = doc.add_paragraph()
    r = p.add_run(f"Overall Score: {_display_value(scorecard.get('jaspen_score'))}")
    r.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = magenta
    if scorecard.get("score_category"):
        p.add_run(f"  ({scorecard.get('score_category')})")

    if scorecard.get("executive_summary"):
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(str(scorecard.get("executive_summary")))

    dims, keys = _scorecard_dim_keys(scorecard)
    if keys:
        doc.add_heading("Scoring", level=1)
        table = doc.add_table(rows=1, cols=3)
        try:
            table.style = "Light Grid Accent 1"
        except Exception:
            pass
        hdr = table.rows[0].cells
        hdr[0].text, hdr[1].text, hdr[2].text = "Criterion", "Score", "Rationale"
        for k in keys:
            d = dims.get(k) if isinstance(dims.get(k), dict) else {}
            cells = table.add_row().cells
            cells[0].text = str(d.get("label") or k)
            cells[1].text = "" if d.get("score") is None else str(d.get("score"))
            cells[2].text = str(d.get("rationale") or "")

    risk_texts = _scorecard_risk_texts(scorecard)
    if risk_texts:
        doc.add_heading("Top Risks", level=1)
        for t in risk_texts:
            doc.add_paragraph(t, style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@export_bp.route("/threads/<thread_id>/scorecard/xlsx", methods=["GET"])
@jwt_required()
def export_scorecard_xlsx(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404
    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404
    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "xlsx")
    if access_error:
        return access_error
    scorecard, error_response = _scorecard_record_for_export(session, thread_id, scorecard_id=request.args.get("scorecard_id"), user_id=user.id)
    if error_response:
        return error_response
    try:
        payload = _xlsx_bytes(scorecard, org=org)
    except RuntimeError as exc:
        return jsonify({"error": str(exc), "code": "xlsx_dependency_missing"}), 503
    filename = f"{_safe_filename_base(scorecard.get('project_name'))}-scorecard.xlsx"
    return _send_bytes(payload, filename=filename, mimetype=XLSX_MIMETYPE)


@export_bp.route("/threads/<thread_id>/scorecard/docx", methods=["GET"])
@jwt_required()
def export_scorecard_docx(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404
    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404
    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "docx")
    if access_error:
        return access_error
    scorecard, error_response = _scorecard_record_for_export(session, thread_id, scorecard_id=request.args.get("scorecard_id"), user_id=user.id)
    if error_response:
        return error_response
    try:
        payload = _docx_bytes(scorecard, org=org)
    except RuntimeError as exc:
        return jsonify({"error": str(exc), "code": "docx_dependency_missing"}), 503
    filename = f"{_safe_filename_base(scorecard.get('project_name'))}-scorecard.docx"
    return _send_bytes(payload, filename=filename, mimetype=DOCX_MIMETYPE)


@export_bp.route("/threads/<thread_id>/scorecard/pdf", methods=["GET"])
@jwt_required()
def export_scorecard_pdf(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "scorecard_pdf")
    if access_error:
        return access_error

    scorecard, error_response = _scorecard_record_for_export(
        session,
        thread_id,
        scorecard_id=request.args.get("scorecard_id"),
        user_id=user.id,
    )
    if error_response:
        return error_response

    payload = _scorecard_pdf_bytes(scorecard, org=org)
    filename = f"{_safe_filename_base(scorecard.get('project_name'))}-scorecard.pdf"
    return _send_bytes(payload, filename=filename, mimetype=PDF_MIMETYPE)


@export_bp.route("/threads/<thread_id>/scorecard/pptx", methods=["GET"])
@jwt_required()
def export_scorecard_pptx(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "scorecard_pptx")
    if access_error:
        return access_error

    scorecard, error_response = _scorecard_record_for_export(
        session,
        thread_id,
        scorecard_id=request.args.get("scorecard_id"),
        user_id=user.id,
    )
    if error_response:
        return error_response

    try:
        payload = _pptx_bytes(scorecard, org=org)
    except RuntimeError as exc:
        return jsonify({"error": str(exc), "code": "pptx_dependency_missing"}), 503

    filename = f"{_safe_filename_base(scorecard.get('project_name'))}-scorecard.pptx"
    return _send_bytes(payload, filename=filename, mimetype=PPTX_MIMETYPE)


@export_bp.route("/threads/<thread_id>/wbs/csv", methods=["GET"])
@jwt_required()
def export_wbs_csv(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    plan_key, _org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "csv")
    if access_error:
        return access_error

    scenario_data = _load_scenarios(user.id) or {}
    thread_data = scenario_data.get(thread_id) if isinstance(scenario_data, dict) else {}
    project_wbs = _resolve_thread_wbs(thread_data, request.args.get("scorecard_id"))
    if not isinstance(project_wbs, dict):
        return jsonify({"error": "No execution plan is available for this thread."}), 404

    payload = _wbs_csv_bytes(project_wbs)
    if payload is None:
        return jsonify({"error": "No WBS tasks are available to export."}), 404

    project_name = session.get("name") or thread_id
    filename = f"{_safe_filename_base(project_name)}-wbs.csv"
    return _send_bytes(payload, filename=filename, mimetype=CSV_MIMETYPE)


@export_bp.route("/threads/<thread_id>/wbs/xlsx", methods=["GET"])
@jwt_required()
def export_wbs_xlsx(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "wbs_xlsx")
    if access_error:
        return access_error

    scenario_data = _load_scenarios(user.id) or {}
    thread_data = scenario_data.get(thread_id) if isinstance(scenario_data, dict) else {}
    project_wbs = _resolve_thread_wbs(thread_data, request.args.get("scorecard_id"))
    if not isinstance(project_wbs, dict):
        return jsonify({"error": "No execution plan is available for this thread."}), 404

    project_name = session.get("name") or project_wbs.get("name") or thread_id
    try:
        payload = _wbs_xlsx_bytes(
            project_wbs,
            project_name=project_name,
            workspace_name=getattr(org, "name", None) or "Personal Workspace",
        )
    except RuntimeError as exc:
        return jsonify({"error": str(exc), "code": "xlsx_dependency_missing"}), 503
    if payload is None:
        return jsonify({"error": "No WBS tasks are available to export."}), 404

    filename = f"{_safe_filename_base(project_name)}-execution-plan.xlsx"
    return _send_bytes(payload, filename=filename, mimetype=XLSX_MIMETYPE)


@export_bp.route("/threads/<thread_id>/conversation/markdown", methods=["GET"])
@jwt_required()
def export_conversation_markdown(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    _plan_key, org, _membership = _resolve_export_context(user, session)
    markdown = _transcript_markdown(session, org=org)
    project_name = session.get("name") or thread_id
    filename = f"{_safe_filename_base(project_name)}-conversation.md"
    return _send_bytes(markdown.encode("utf-8"), filename=filename, mimetype=MARKDOWN_MIMETYPE)


@export_bp.route("/threads/<thread_id>/conversation/pdf", methods=["GET"])
@jwt_required()
def export_conversation_pdf(thread_id):
    user = User.query.filter_by(id=str(get_jwt_identity())).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    sessions = load_user_sessions(user.id) or {}
    session = _resolve_thread_session(sessions, thread_id)
    if not isinstance(session, dict):
        return jsonify({"error": "Thread not found"}), 404

    plan_key, org, _membership = _resolve_export_context(user, session)
    access_error = _require_export_plan(plan_key, "pdf")
    if access_error:
        return access_error

    markdown = _transcript_markdown(session, org=org)
    project_name = session.get("name") or thread_id
    payload = _markdown_to_pdf_bytes(project_name, markdown)
    filename = f"{_safe_filename_base(project_name)}-conversation.pdf"
    return _send_bytes(payload, filename=filename, mimetype=PDF_MIMETYPE)
